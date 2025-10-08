import os
from PIL import Image
import pdfkit
from reportlab.pdfgen import canvas
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4



def convert_to_pdf(input_file):
    # Pega a extensão do arquivo (exemplo: .png, .html, etc)
    ext = os.path.splitext(input_file)[1].lower()

    # Cria o nome do arquivo de saída (mesmo nome mas .pdf)
    output_file = input_file.replace(ext, ".pdf")

    # Se for imagem -> usa Pillow
    if ext in [".jpg", ".jpeg", ".png"]:
        print("Convertendo imagem para PDF...")
        img = Image.open(input_file).convert("RGB")
        img.save(output_file)

    # Se for HTML -> usa pdfkit (precisa do wkhtmltopdf instalado)
    elif ext == ".html":
        print("Convertendo HTML para PDF...")
        pdfkit.from_file(input_file, output_file)
        # Se for PowerPoint (.pptx) -> usa python-pptx + reportlab
    elif ext in [".ppt", ".pptx"]:
        print("Convertendo PowerPoint para PDF...")


        prs = Presentation(input_file)
        output_file = input_file.replace(ext, ".pdf")

        c = canvas.Canvas(output_file, pagesize=A4)
        width, height = A4

        for slide_num, slide in enumerate(prs.slides, start=1):
            y = height - 50
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, y, f"Slide {slide_num}")
            y -= 30

            # percorre todos os shapes (caixas de texto, títulos etc)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto = shape.text.strip()
                    if texto:
                        c.setFont("Helvetica", 11)
                        for linha in texto.splitlines():
                            c.drawString(60, y, linha[:100])  # corta se for muito longo
                            y -= 15
                            if y < 50:  # cria nova página se acabar o espaço
                                c.showPage()
                                y = height - 50

            c.showPage()

        c.save()


    # Se for TXT simples -> converte "na unha"
    elif ext == ".txt":
        print("Convertendo TXT para PDF...")
        c = canvas.Canvas(output_file)
        with open(input_file, "r", encoding="utf-8") as f:
            linhas = f.readlines()
            y = 800  # posição inicial no PDF
            for linha in linhas:
                c.drawString(50, y, linha.strip())  # escreve cada linha
                y -= 20  # sobe uma linha
        c.save()

    else:
        print("Formato não suportado:", ext)
        return None

    print("✅ Arquivo convertido com sucesso:", output_file)
    return output_file


# Programa principal
if __name__ == "__main__":
    print("=== Conversor Simples para PDF ===")
    caminho = input("Digite o caminho completo do arquivo que deseja converter: ")

    # Verifica se o arquivo existe
    if os.path.exists(caminho):
        convert_to_pdf(caminho)
    else:
        print("❌ Arquivo não encontrado. Verifique o caminho e tente novamente.")
